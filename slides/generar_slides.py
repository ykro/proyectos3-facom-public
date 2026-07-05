#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generador de slides (PPTX) para el curso
"Proyectos III — Diseño de aplicaciones móviles" (Universidad Galileo).

Tema CLARO (paleta vino del curso). Tipografía grande para aprovechar el
espacio. Un deck por sesión de contenido: 1, 2, 3, 4, 6, 7, 8, 9.
(Las sesiones 5 y 10 son exámenes y no llevan deck.)

Cada deck incluye, según la sesión: portada, agenda, slides de contenido,
diagramas de flujo, estadísticas con fuente citada, capturas reales de las
apps del curso (assets/), laboratorio y cierre.

Cada sesión de contenido incluye además un slide de "Demo del docente":
la mini-app distinta que el docente construye EN VIVO para ilustrar el
concepto del día (no es lo que el estudiante entrega).

Ejecutar con:
    uv run --with python-pptx python slides/generar_slides.py

Produce: slides/sesion-01.pptx ... sesion-09.pptx (8 archivos).
Requiere las capturas en slides/assets/ (ver slides/assets/README si falta).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paleta CLARA (vino) — la misma paleta que usan las apps del curso
# ---------------------------------------------------------------------------
CREMA     = RGBColor(0xFD, 0xF6, 0xEC)   # fondo principal
CREMA_2   = RGBColor(0xF6, 0xEC, 0xDD)   # paneles suaves
CARD      = RGBColor(0xFF, 0xFF, 0xFF)   # tarjetas
VINO      = RGBColor(0x6D, 0x2E, 0x46)   # primario
TERRACOTA = RGBColor(0xC8, 0x4B, 0x31)   # acento cálido
MOSTAZA   = RGBColor(0xE3, 0xB2, 0x3C)   # acento
VERDE     = RGBColor(0x4F, 0x77, 0x2D)   # extra
AZUL      = RGBColor(0x3A, 0x6B, 0x7E)   # extra
TEXTO     = RGBColor(0x2B, 0x2B, 0x2B)   # texto principal
MUTED     = RGBColor(0x6F, 0x64, 0x5B)   # texto secundario (AA sobre crema; igual que el sistema de diseño)
LINE      = RGBColor(0xEA, 0xDD, 0xC8)   # separadores

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COURSE = "Proyectos III — Diseño de aplicaciones móviles"
SUBTITLE_COURSE = "Licenciatura en Comunicación y Diseño · Universidad Galileo · Docente: Adrián Catalán"

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
PHONE_RATIO = 844 / 390  # alto / ancho de las capturas


# ---------------------------------------------------------------------------
# Utilidades de bajo nivel
# ---------------------------------------------------------------------------
def _set_bg(slide, color=CREMA):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color, rounded=False, line_color=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_text(slide, left, top, width, height, runs, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, space_after=6):
    """runs: lista de (texto, tamaño_pt, color, negrita, nivel)."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, color, bold, level) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.level = level
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Calibri"
    return box


def _add_bullets(slide, left, top, width, height, bullets, size=22,
                 color=TEXTO, bullet_char="•  ", space=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        text, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(space)
        run = p.add_run()
        prefix = bullet_char if lvl == 0 else "–  "
        run.text = prefix + text
        run.font.size = Pt(size if lvl == 0 else size - 3)
        run.font.color.rgb = color if lvl == 0 else MUTED
        run.font.name = "Calibri"
    return box


def _title_block(slide, title, kicker=None, bar=TERRACOTA):
    """Encabezado estándar: barra superior + kicker + título + subrayado."""
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.16), bar)
    if kicker:
        _add_text(slide, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.5),
                  [(kicker.upper(), 18, TERRACOTA, True, 0)])
    _add_text(slide, Inches(0.9), Inches(1.0), Inches(11.6), Inches(1.1),
              [(title, 38, VINO, True, 0)])
    _add_rect(slide, Inches(0.95), Inches(1.95), Inches(1.7), Inches(0.07), MOSTAZA)


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, CREMA)
    return slide


# ---------------------------------------------------------------------------
# Tipos de diapositiva
# ---------------------------------------------------------------------------
def add_cover(prs, session_no, title):
    slide = _blank(prs)
    # banda lateral vino
    _add_rect(slide, 0, 0, Inches(0.45), SLIDE_H, VINO)
    _add_rect(slide, Inches(0.45), 0, Inches(0.12), SLIDE_H, MOSTAZA)
    _add_text(slide, Inches(1.1), Inches(1.35), Inches(11.4), Inches(0.7),
              [(f"SESIÓN {session_no:02d}", 26, TERRACOTA, True, 0)])
    _add_text(slide, Inches(1.05), Inches(2.05), Inches(11.6), Inches(2.6),
              [(title, 52, VINO, True, 0)])
    _add_rect(slide, Inches(1.12), Inches(4.75), Inches(2.4), Inches(0.09), MOSTAZA)
    _add_text(slide, Inches(1.1), Inches(5.1), Inches(11.4), Inches(0.6),
              [(COURSE, 23, TEXTO, True, 0)])
    _add_text(slide, Inches(1.1), Inches(5.7), Inches(11.4), Inches(0.8),
              [(SUBTITLE_COURSE, 15, MUTED, False, 0)])
    return slide


def add_agenda(prs, items):
    slide = _blank(prs)
    _title_block(slide, "Agenda", kicker="Sesión de hoy")
    _add_bullets(slide, Inches(0.98), Inches(2.35), Inches(11.4), Inches(4.7),
                 items, size=26, bullet_char="●  ", space=16)
    return slide


def add_content(prs, title, bullets, kicker=None):
    slide = _blank(prs)
    _title_block(slide, title, kicker=kicker)
    _add_bullets(slide, Inches(0.98), Inches(2.3), Inches(11.5), Inches(4.9),
                 bullets, size=23, space=15)
    return slide


def add_stats(prs, title, cards, kicker="Con datos"):
    """cards: lista de dict {num, label, source}. Rejilla 2x2 (o 1xN)."""
    slide = _blank(prs)
    _title_block(slide, title, kicker=kicker, bar=AZUL)
    n = len(cards)
    cols = 2 if n > 2 else n
    rows = (n + cols - 1) // cols
    gap = Inches(0.38)
    area_left, area_top = Inches(0.95), Inches(2.25)
    area_w = Inches(11.45)
    cw = (area_w - gap * (cols - 1)) / cols
    ch = Inches(2.28) if rows > 1 else Inches(3.7)
    accents = [TERRACOTA, VINO, VERDE, AZUL, MOSTAZA, TERRACOTA]
    for i, c in enumerate(cards):
        r, col = divmod(i, cols)
        left = area_left + col * (cw + gap)
        top = area_top + r * (ch + gap)
        _add_rect(slide, left, top, cw, ch, CARD, rounded=True, line_color=LINE)
        _add_rect(slide, left, top, Inches(0.13), ch, accents[i % len(accents)])
        num_sz = 44 if rows > 1 else 54
        _add_text(slide, left + Inches(0.38), top + Inches(0.16), cw - Inches(0.65), Inches(0.9),
                  [(c["num"], num_sz, accents[i % len(accents)], True, 0)])
        # etiqueta entre el número y la fuente (fuente anclada abajo)
        _add_text(slide, left + Inches(0.38), top + Inches(1.02), cw - Inches(0.65),
                  ch - Inches(1.5),
                  [(c["label"], 16, TEXTO, True, 0)])
        _add_text(slide, left + Inches(0.38), top + ch - Inches(0.42), cw - Inches(0.65), Inches(0.36),
                  [("Fuente: " + c["source"], 10.5, MUTED, False, 0)])
    return slide


def add_diagram(prs, title, steps, kicker="Cómo funciona", footnote=None, bar=VERDE):
    """steps: lista de (label) o (label, sublabel). Fila horizontal con flechas."""
    slide = _blank(prs)
    _title_block(slide, title, kicker=kicker, bar=bar)
    n = len(steps)
    total_w = Inches(11.6)
    arrow_w = Inches(0.55)
    box_w = (total_w - arrow_w * (n - 1)) / n
    box_h = Inches(2.1)
    left0 = Inches(0.9)
    top = Inches(3.3)
    palette = [VINO, TERRACOTA, VERDE, AZUL, MOSTAZA]
    x = left0
    for i, step in enumerate(steps):
        label, sub = step if isinstance(step, tuple) else (step, None)
        color = palette[i % len(palette)]
        _add_rect(slide, x, top, box_w, box_h, CARD, rounded=True, line_color=LINE)
        _add_rect(slide, x, top, box_w, Inches(0.14), color)
        _add_text(slide, x + Inches(0.15), top + Inches(0.5), box_w - Inches(0.3),
                  box_h - Inches(0.7),
                  [(label, 19, VINO, True, 0)] + ([(sub, 14, MUTED, False, 0)] if sub else []),
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += box_w
        if i < n - 1:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(0.06),
                                        top + box_h/2 - Inches(0.22), arrow_w - Inches(0.12), Inches(0.44))
            ar.fill.solid(); ar.fill.fore_color.rgb = MOSTAZA
            ar.line.fill.background(); ar.shadow.inherit = False
            x += arrow_w
    if footnote:
        _add_text(slide, Inches(0.95), Inches(5.9), Inches(11.4), Inches(0.8),
                  [(footnote, 17, TEXTO, False, 0)])
    return slide


def add_compare(prs, title, left_title, left_items, right_title, right_items,
                kicker="Comparación", left_color=AZUL, right_color=VINO):
    slide = _blank(prs)
    _title_block(slide, title, kicker=kicker, bar=MOSTAZA)
    top = Inches(2.35)
    h = Inches(4.6)
    w = Inches(5.55)
    gap = Inches(0.35)
    left0 = Inches(0.95)
    for idx, (ttl, items, color, lx) in enumerate([
            (left_title, left_items, left_color, left0),
            (right_title, right_items, right_color, left0 + w + gap)]):
        _add_rect(slide, lx, top, w, h, CARD, rounded=True, line_color=LINE)
        _add_rect(slide, lx, top, w, Inches(0.85), color, rounded=True)
        _add_rect(slide, lx, top + Inches(0.5), w, Inches(0.35), color)
        _add_text(slide, lx + Inches(0.35), top + Inches(0.12), w - Inches(0.6), Inches(0.7),
                  [(ttl, 24, CARD, True, 0)], anchor=MSO_ANCHOR.MIDDLE)
        _add_bullets(slide, lx + Inches(0.35), top + Inches(1.1), w - Inches(0.7), h - Inches(1.3),
                     items, size=19, space=12)
    return slide


def add_screens(prs, title, images, kicker="En pantalla", bullets=None):
    """images: lista de (nombre_asset, caption). 1 o 2 teléfonos."""
    slide = _blank(prs)
    _title_block(slide, title, kicker=kicker, bar=TERRACOTA)
    imgs = images[:2]
    n = len(imgs)
    if n == 1 and bullets:
        # 1 teléfono a la derecha + bullets a la izquierda
        ph_h = Inches(4.6)
        ph_w = ph_h / PHONE_RATIO
        left = Inches(10.05)
        top = Inches(2.05)
        _place_phone(slide, imgs[0], left, top, ph_h, ph_w)
        _add_bullets(slide, Inches(0.98), Inches(2.5), Inches(8.4), Inches(4.4),
                     bullets, size=22, space=15)
    else:
        ph_h = Inches(4.45)
        ph_w = ph_h / PHONE_RATIO
        gap = Inches(1.1)
        total = ph_w * n + gap * (n - 1)
        start = (SLIDE_W - total) / 2
        top = Inches(2.2)
        for i, im in enumerate(imgs):
            left = start + i * (ph_w + gap)
            _place_phone(slide, im, left, top, ph_h, ph_w)
    return slide


def _place_phone(slide, image, left, top, ph_h, ph_w):
    name, caption = image
    path = os.path.join(ASSETS, name + ".png")
    # marco sutil
    frame = _add_rect(slide, left - Inches(0.06), top - Inches(0.06),
                      ph_w + Inches(0.12), ph_h + Inches(0.12), CARD,
                      rounded=True, line_color=VINO)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, height=ph_h)
    else:
        _add_text(slide, left, top + ph_h/2 - Inches(0.4), ph_w, Inches(0.8),
                  [("[captura: " + name + "]", 12, MUTED, False, 0)],
                  align=PP_ALIGN.CENTER)
    if caption:
        _add_text(slide, left - Inches(0.9), top + ph_h + Inches(0.14), ph_w + Inches(1.8),
                  Inches(0.5), [(caption, 13, VINO, True, 0)], align=PP_ALIGN.CENTER)


def _png_size(path):
    """Lee ancho/alto de un PNG sin dependencias (cabecera IHDR)."""
    try:
        with open(path, "rb") as f:
            head = f.read(24)
    except OSError:
        return None
    if len(head) < 24 or head[:8] != b"\x89PNG\r\n\x1a\n":
        return None
    import struct
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def _place_image_fit(slide, path, box_left, box_top, box_w, box_h, frame=True):
    """Coloca una imagen ajustada DENTRO de la caja, respetando su proporción
    (sirve para capturas verticales de teléfono y horizontales de navegador)."""
    size = _png_size(path) if os.path.exists(path) else None
    if not size:
        _add_rect(slide, box_left, box_top, box_w, box_h, CREMA_2, rounded=True, line_color=LINE)
        _add_text(slide, box_left, box_top + box_h // 2 - Inches(0.3), box_w, Inches(0.6),
                  [("[captura pendiente]", 13, MUTED, False, 0)], align=PP_ALIGN.CENTER)
        return
    w_px, h_px = size
    ratio = h_px / w_px  # alto / ancho
    disp_w = box_w
    disp_h = Emu(int(disp_w * ratio))
    if disp_h > box_h:
        disp_h = box_h
        disp_w = Emu(int(disp_h / ratio))
    left = box_left + (box_w - disp_w) // 2
    top = box_top + (box_h - disp_h) // 2
    if frame:
        _add_rect(slide, left - Inches(0.05), top - Inches(0.05),
                  disp_w + Inches(0.1), disp_h + Inches(0.1), CARD,
                  rounded=True, line_color=VINO)
    slide.shapes.add_picture(path, left, top, width=disp_w, height=disp_h)


def add_pair(prs, title, mockup, mockup_caption, real, real_caption,
             kicker="Meta y resultado"):
    """Empareja el mock-up de referencia (\"esto vamos a construir\") con la
    captura real (\"así quedó, generado por la IA\"), lado a lado."""
    slide = _blank(prs)
    _title_block(slide, title, kicker=kicker, bar=MOSTAZA)
    cols = [
        (mockup, mockup_caption, "1 · Esto vamos a construir", "Mock-up de referencia", AZUL),
        (real,   real_caption,   "2 · Así quedó, generado por la IA", "Captura real de la app", VINO),
    ]
    w = Inches(5.85)
    gap = Inches(0.4)
    left0 = Inches(0.72)
    top = Inches(2.2)
    header_h = Inches(0.72)
    for i, (name, cap, label, sub, color) in enumerate(cols):
        lx = left0 + i * (w + gap)
        _add_rect(slide, lx, top, w, header_h, color, rounded=True)
        _add_text(slide, lx + Inches(0.28), top + Inches(0.05), w - Inches(0.55),
                  header_h - Inches(0.1),
                  [(label, 17, CARD, True, 0), (sub, 12, CREMA, False, 0)],
                  anchor=MSO_ANCHOR.MIDDLE, space_after=1)
        box_top = top + header_h + Inches(0.18)
        box_h = Inches(3.7)
        _place_image_fit(slide, os.path.join(ASSETS, name + ".png"), lx, box_top, w, box_h)
        if cap:
            _add_text(slide, lx, box_top + box_h + Inches(0.05), w, Inches(0.5),
                      [(cap, 13, VINO, True, 0)], align=PP_ALIGN.CENTER)
    return slide


def add_demo(prs, app_name, concept, image, kicker="Demo del docente · en vivo"):
    """Slide de la mini-app que el DOCENTE construye en vivo esta sesión para
    ilustrar el concepto del día. NO es lo que el estudiante entrega: el
    laboratorio del estudiante sigue en la app insignia + su propia idea.
    concept: lista de líneas cortas (1–2). image: nombre de asset (sin .png)."""
    slide = _blank(prs)
    _title_block(slide, app_name, kicker=kicker, bar=VERDE)
    # Columna izquierda: concepto + aclaración
    txt_left = Inches(0.98)
    txt_w = Inches(7.2)
    _add_rect(slide, txt_left, Inches(2.35), txt_w, Inches(0.62), CREMA_2, rounded=True)
    _add_text(slide, txt_left + Inches(0.3), Inches(2.45), txt_w - Inches(0.6), Inches(0.5),
              [("El docente la construye en vivo — no es tu entrega", 17, VINO, True, 0)],
              anchor=MSO_ANCHOR.MIDDLE)
    _add_bullets(slide, txt_left, Inches(3.25), txt_w, Inches(2.3),
                 concept, size=22, space=14)
    _add_text(slide, txt_left, Inches(5.75), txt_w, Inches(1.3),
              [("Ilustra el concepto del día con una app distinta. Tu laboratorio "
                "sigue en la app insignia y en tu propia idea.", 15, MUTED, False, 0)])
    # Columna derecha: captura real de la mini-app
    box_left = Inches(8.55)
    box_w = Inches(4.0)
    _place_image_fit(slide, os.path.join(ASSETS, image + ".png"),
                     box_left, Inches(2.3), box_w, Inches(4.7))
    _add_text(slide, box_left, Inches(7.02), box_w, Inches(0.4),
              [("Captura real de la demo", 12, VINO, True, 0)], align=PP_ALIGN.CENTER)
    return slide


def add_glossary(prs, title, terms, kicker="Vocabulario mínimo"):
    """terms: lista de (término, definición corta). Rejilla de 2 columnas."""
    slide = _blank(prs)
    _title_block(slide, title, kicker=kicker, bar=VERDE)
    n = len(terms)
    per = (n + 1) // 2
    col_w = Inches(5.85)
    gap = Inches(0.4)
    left0 = Inches(0.72)
    top0 = Inches(2.2)
    row_h = Inches(0.86)
    row_gap = Inches(0.11)
    for i, (term, definition) in enumerate(terms):
        col = i // per
        row = i % per
        lx = left0 + col * (col_w + gap)
        ty = top0 + row * (row_h + row_gap)
        _add_rect(slide, lx, ty, col_w, row_h, CARD, rounded=True, line_color=LINE)
        _add_rect(slide, lx, ty, Inches(0.1), row_h, MOSTAZA)
        _add_text(slide, lx + Inches(0.28), ty + Inches(0.06), col_w - Inches(0.5),
                  row_h - Inches(0.12),
                  [(term, 15, VINO, True, 0), (definition, 12.5, TEXTO, False, 0)],
                  anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    return slide


def add_lab(prs, title_lab, bullets):
    slide = _blank(prs)
    _add_rect(slide, Inches(0.7), Inches(0.7), Inches(11.93), Inches(6.1), CARD,
              rounded=True, line_color=LINE)
    _add_rect(slide, Inches(0.7), Inches(0.7), Inches(0.22), Inches(6.1), MOSTAZA)
    _add_text(slide, Inches(1.15), Inches(1.05), Inches(11.0), Inches(0.5),
              [("LABORATORIO DE HOY", 20, TERRACOTA, True, 0)])
    _add_text(slide, Inches(1.15), Inches(1.6), Inches(11.0), Inches(1.0),
              [(title_lab, 34, VINO, True, 0)])
    _add_bullets(slide, Inches(1.2), Inches(2.85), Inches(10.8), Inches(3.7),
                 bullets, size=23, space=15)
    return slide


def add_closing(prs, recap, tareas):
    slide = _blank(prs)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.16), VINO)
    _add_text(slide, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.5),
              [("CIERRE", 18, TERRACOTA, True, 0)])
    _add_text(slide, Inches(0.9), Inches(1.0), Inches(11.5), Inches(0.9),
              [("Lo que nos llevamos y la tarea", 36, VINO, True, 0)])
    _add_rect(slide, Inches(0.95), Inches(1.95), Inches(1.7), Inches(0.07), MOSTAZA)
    # columna izquierda: recap
    _add_rect(slide, Inches(0.9), Inches(2.35), Inches(5.75), Inches(4.55), CREMA_2, rounded=True)
    _add_text(slide, Inches(1.2), Inches(2.6), Inches(5.2), Inches(0.5),
              [("En resumen", 22, VERDE, True, 0)])
    _add_bullets(slide, Inches(1.2), Inches(3.2), Inches(5.2), Inches(3.6),
                 recap, size=20, space=13)
    # columna derecha: tareas
    _add_rect(slide, Inches(6.85), Inches(2.35), Inches(5.6), Inches(4.55), CREMA_2, rounded=True)
    _add_text(slide, Inches(7.15), Inches(2.6), Inches(5.0), Inches(0.5),
              [("Para la próxima", 22, TERRACOTA, True, 0)])
    _add_bullets(slide, Inches(7.15), Inches(3.2), Inches(5.0), Inches(3.6),
                 tareas, size=20, space=13)
    return slide


# ---------------------------------------------------------------------------
# Construcción de un deck a partir de la especificación
# ---------------------------------------------------------------------------
def build_deck(spec, out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover(prs, spec["session"], spec["title"])
    add_agenda(prs, spec["agenda"])
    for c in spec["content"]:
        add_content(prs, c["title"], c["bullets"], kicker=c.get("kicker"))
    if spec.get("demo"):
        d = spec["demo"]
        add_demo(prs, d["app"], d["concept"], d["image"],
                 kicker=d.get("kicker", "Demo del docente · en vivo"))
    for g in spec.get("glossary", []):
        add_glossary(prs, g["title"], g["terms"], kicker=g.get("kicker", "Vocabulario mínimo"))
    for d in spec.get("diagrams", []):
        add_diagram(prs, d["title"], d["steps"], kicker=d.get("kicker", "Cómo funciona"),
                    footnote=d.get("footnote"))
    for cmp in spec.get("compare", []):
        add_compare(prs, cmp["title"], cmp["left_title"], cmp["left"],
                    cmp["right_title"], cmp["right"], kicker=cmp.get("kicker", "Comparación"))
    for s in spec.get("stats", []):
        add_stats(prs, s["title"], s["cards"], kicker=s.get("kicker", "Con datos"))
    for sc in spec.get("screens", []):
        add_screens(prs, sc["title"], sc["images"], kicker=sc.get("kicker", "En pantalla"),
                    bullets=sc.get("bullets"))
    for p in spec.get("pairs", []):
        add_pair(prs, p["title"], p["mockup"], p.get("mockup_caption"),
                 p["real"], p.get("real_caption"), kicker=p.get("kicker", "Meta y resultado"))
    add_lab(prs, spec["lab"]["title"], spec["lab"]["bullets"])
    add_closing(prs, spec["closing"]["recap"], spec["closing"]["tareas"])

    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# Fuentes de las estadísticas (citadas en cada slide)
# ---------------------------------------------------------------------------
# DataReportal (We Are Social & Meltwater), Digital 2025 Global / Guatemala
# Stack Overflow Developer Survey 2025 · Figma AI Design Report 2024
# Sensor Tower State of Mobile 2024/2025 · Statista Google Play 2024
# LinkedIn Future of Work: AI at Work 2025 · Lindgaard et al. 2006 (50 ms)

# ---------------------------------------------------------------------------
# Contenido de cada sesión
# ---------------------------------------------------------------------------
DECKS = [
    {
        "session": 1,
        "title": "Introducción, UX/UI móvil y el diseñador-creador con IA",
        "agenda": [
            "Bienvenida y de qué trata el curso",
            "El diseñador-creador con IA",
            "Panorama de herramientas",
            "Principios de UX/UI móvil",
            "Los datos: por qué móvil y por qué ahora",
            "Demo en vivo y Laboratorio 1",
        ],
        "content": [
            {"title": "Desarrollo de software asistido por IA", "kicker": "La idea central", "bullets": [
                "Construimos apps móviles reales sin escribir código",
                "Dirigimos a la IA en español con instrucciones claras",
                "Tu fortaleza —diseño, UX, identidad visual, narrativa— guía a la IA",
                "Siempre decimos \"desarrollo asistido por IA\", nunca \"vibe coding\"",
            ]},
            {"title": "Dos proyectos a lo largo del trimestre", "kicker": "Hacia dónde vamos", "bullets": [
                "Proyecto 1 (PWA): \"Ruta del Café\", semanas 1–3",
                "Se evalúa en el Examen Parcial (sesión 5)",
                "Proyecto Final (nativa): \"Bitácora Visual\", semanas 4, 6 y 7",
                "Se entrega y se muestra en el Demo Day (sesión 10)",
            ]},
            {"title": "Panorama de herramientas", "kicker": "Con qué construimos", "bullets": [
                "Google AI Studio (modo Build): todo desde el navegador (Chrome)",
                "La IA de AI Studio genera el código; modelo por defecto Gemini 3.5 Flash",
                "PWA (React) y app nativa Android (Kotlin + Compose), en la misma herramienta",
                "Firebase: Firestore, Authentication, Storage",
                "Cloud Run para publicar · AI Chips (Nano Banana)",
                "Todo gratis; solo necesitas una cuenta de Google",
            ]},
            {"title": "Principios de UX/UI móvil", "kicker": "Diseñar para el teléfono", "bullets": [
                "Mobile-first: pantalla pequeña, una columna, jerarquía clara",
                "Áreas táctiles grandes (mínimo ~44px) y buen contraste",
                "Diseñar para el pulgar: lo importante al alcance",
                "Patrones: tarjetas (cards), filtros, vista de detalle, navegación",
                "Cuidar los estados: vacío, cargando y error",
            ]},
            {"title": "GitHub: guardar y compartir tu trabajo", "kicker": "Opcional y muy simple", "bullets": [
                "Un repositorio (repo) es como una carpeta en la nube con el historial de tu proyecto",
                "Sirve para guardar versiones, no perder trabajo y compartir con otras personas",
                "AI Studio puede exportar tu app a GitHub con un clic (opcional en el curso)",
                "Piénsalo como tu portafolio de código: tu app vive ahí y la puedes mostrar",
            ]},
        ],
        "glossary": [
            {"title": "Glosario mínimo de IA generativa", "kicker": "Vocabulario para diseño", "terms": [
                ("Prompt", "La instrucción en español que le das a la IA para crear o cambiar algo."),
                ("Modelo (Gemini 3.5 Flash)", "El \"cerebro\" de IA que genera; el que usa AI Studio por defecto."),
                ("Token", "La unidad mínima de texto que la IA procesa (trozos de palabras)."),
                ("Generar", "Cuando la IA produce el código y la app a partir de tu prompt."),
                ("Iterar / refinar", "Mejorar la app pidiendo un cambio a la vez."),
                ("Checkpoint", "Un punto guardado al que puedes volver (como guardar la partida)."),
                ("Alucinación", "Cuando la IA inventa algo incorrecto; siempre hay que verificar."),
                ("AI Chips (Nano Banana)", "Servicios de Google dentro de AI Studio; Nano Banana genera imágenes."),
                ("Firebase", "La nube de Google: datos (Firestore), login (Auth) y archivos (Storage)."),
            ]},
        ],
        "diagrams": [
            {"title": "¿Cómo hace la IA todo esto?", "kicker": "Para diseño",
             "steps": [("Prompt en español", "tu instrucción"),
                       ("La IA escribe el código", "interfaz + lógica + datos"),
                       ("Vista previa", "la app corriendo")],
             "footnote": "Tú diriges con criterio de diseño; la IA es como un desarrollador junior muy "
                         "rápido: ella escribe el código, tú decides qué se construye y cómo se ve."},
            {"title": "Cómo se dirige a la IA", "kicker": "Flujo de trabajo",
             "steps": [("Prompt inicial", "detallado"), ("Vista previa", "revisar"),
                       ("Refinar", "un cambio"), ("Checkpoint", "guardar")],
             "footnote": "Es un ciclo: revisas la vista previa, pides un cambio a la vez y guardas checkpoints. "
                         "View diff te muestra qué cambió; regla de los 2 strikes si algo se atora."},
        ],
        "stats": [
            {"title": "Por qué móvil y por qué ahora", "kicker": "Los datos", "cards": [
                {"num": "51.6%", "label": "del tráfico web mundial ya viene de teléfonos móviles",
                 "source": "DataReportal, Digital 2025 Global Overview"},
                {"num": "20.4 M", "label": "conexiones móviles en Guatemala (≈110% de la población)",
                 "source": "DataReportal, Digital 2025 Guatemala"},
                {"num": "84%", "label": "de desarrolladores ya usa o planea usar IA en su trabajo",
                 "source": "Stack Overflow Developer Survey 2025"},
                {"num": "59%", "label": "de diseñadores y desarrolladores ya usa IA en su flujo",
                 "source": "Figma, AI Design Report 2024"},
            ]},
        ],
        "screens": [
            {"title": "Tu primera app: Ruta del Café", "kicker": "A dónde llegamos hoy",
             "images": [("ruta-home", "Pantalla principal: tarjetas y filtro")],
             "bullets": [
                 "Con un solo prompt inicial detallado obtienes una app así",
                 "Encabezado, filtro por categoría y tarjetas (cards)",
                 "Cada tarjeta: imagen, nombre, zona y categoría",
                 "Hoy la construyes y aplicas tu primer refinamiento",
             ]},
        ],
        "pairs": [
            {"title": "Ruta del Café: la meta y el resultado",
             "mockup": "ruta-home", "mockup_caption": "El mock-up de referencia (la guía)",
             "real": "ruta-real", "real_caption": "La app generada por la IA en AI Studio"},
        ],
        "demo": {
            "app": "Antojitos Chapines",
            "image": "demo-antojitos",
            "concept": [
                "Catálogo de antojitos callejeros de Guatemala (PWA)",
                "Mismo patrón que Ruta del Café: tarjetas, filtro y estética de marca",
                "El docente la arma en vivo con un solo prompt inicial detallado",
            ],
        },
        "lab": {"title": "Laboratorio 1 — Tu primera app en AI Studio", "bullets": [
            "Crea una app nueva en AI Studio (modo Build)",
            "Pega el prompt inicial de \"Ruta del Café\"",
            "Verifica el encabezado y las tarjetas de cafeterías",
            "Crea un checkpoint y aplica el refinamiento de colores",
            "Usa View diff para ver el cambio",
        ]},
        "closing": {
            "recap": [
                "Construimos software dirigiendo a la IA en español",
                "Un buen prompt inicial ahorra muchos refinamientos",
                "Checkpoints, View diff y la regla de los 2 strikes",
            ],
            "tareas": [
                "Tarea 1: investigación sobre principios de UI/UX móvil",
                "Laboratorio 1: tu primera app con un refinamiento aplicado",
                "Guarda tus prompts en un documento desde hoy",
            ],
        },
    },
    {
        "session": 2,
        "title": "De idea a interfaz y prototipo",
        "agenda": [
            "Repaso de la sesión 1",
            "Prototipos navegables",
            "Identidad visual como sistema",
            "AI Chips y Nano Banana",
            "Modo de anotación",
            "Laboratorio 2 y tarea",
        ],
        "content": [
            {"title": "Del mockup al prototipo navegable", "kicker": "Concepto", "bullets": [
                "Un prototipo se puede tocar: pantallas conectadas de verdad",
                "Flujo lista → detalle → volver",
                "Incluye estados y microinteracciones",
                "Es más que una imagen estática: es una experiencia",
            ]},
            {"title": "Identidad visual como sistema", "kicker": "Tu marca", "bullets": [
                "Nombre, paleta de color, tipografía y tono",
                "Colores en hexadecimal: #6D2E46, no \"vino\"",
                "Consistencia en toda la app",
                "Contraste accesible: tu criterio de diseño marca la diferencia",
            ]},
            {"title": "AI Chips: servicios de Google sin configuración", "kicker": "Herramientas", "bullets": [
                "Se agregan dentro de AI Studio sin instalar nada",
                "Nano Banana: generar y editar imágenes",
                "Útil para portadas, banners e ilustraciones",
                "Las imágenes pueden alimentar el campo imagenUrl",
            ]},
            {"title": "Modo de anotación", "kicker": "Cambios precisos", "bullets": [
                "Haz clic sobre el elemento en la vista previa",
                "Describe el ajuste justo ahí",
                "Más preciso que describirlo solo con palabras",
                "Ideal para pulir detalles visuales",
            ]},
            {"title": "Dónde viven los colores y estilos (PWA)", "kicker": "Vocabulario de diseño, no programación", "bullets": [
                "La pestaña Code muestra los archivos que generó la IA; no hay que escribirlos, solo saber leerlos",
                "Los colores suelen vivir juntos (config de estilos): ahí cambias la paleta de toda la app",
                "La tipografía y los tamaños se definen una vez y se reutilizan: consistencia por diseño",
                "Puedes ajustar a mano un valor (por ejemplo un hexadecimal) y ver el cambio en la vista previa",
                "Aun así, lo normal es pedir el cambio por prompt; editar a mano es solo para retoques finos",
            ]},
        ],
        "diagrams": [
            {"title": "La navegación de un prototipo", "kicker": "Lista → detalle → volver",
             "steps": [("Lista", "tarjetas"), ("Tocar", "una tarjeta"),
                       ("Detalle", "info completa"), ("Volver", "a la lista")],
             "footnote": "Este mismo patrón lo usarás en casi todas las apps: una pantalla que muestra "
                         "muchos elementos y otra que muestra uno a fondo."},
        ],
        "screens": [
            {"title": "Lista y detalle en Ruta del Café", "kicker": "Lo que agregamos hoy",
             "images": [("ruta-home", "Lista con filtro"), ("ruta-detalle", "Vista de detalle")]},
        ],
        "demo": {
            "app": "Adopta un Peludo",
            "image": "demo-adopta",
            "concept": [
                "Mascotas en adopción, con su ficha de detalle (PWA)",
                "Ilustra el flujo lista → detalle → volver y una identidad visual propia",
                "El docente la construye y la pule en vivo con el modo de anotación",
            ],
        },
        "lab": {"title": "Laboratorio 2 — Prototipo navegable", "bullets": [
            "Guarda un checkpoint antes de refinar",
            "Agrega el filtro por categoría (con \"Todas\")",
            "Agrega la vista de detalle con botón \"Volver\"",
            "Aplica tu propia identidad visual (Tarea 2)",
            "Pule un detalle con el modo de anotación",
        ]},
        "closing": {
            "recap": [
                "Un prototipo navegable se puede tocar de verdad",
                "La identidad visual es un sistema, no un color suelto",
                "Nano Banana para imágenes; anotación para precisión",
            ],
            "tareas": [
                "Tarea 2: define la identidad visual de tu app",
                "Laboratorio 2: prototipo navegable con identidad propia",
                "Reúne referencias de paletas y tipografías",
            ],
        },
    },
    {
        "session": 3,
        "title": "Datos y app instalable (PWA)",
        "agenda": [
            "Repaso de la sesión 2",
            "Firestore: datos en la nube",
            "Tiempo real y CRUD",
            "Auto-provisioning de Firebase",
            "Qué hace instalable a una PWA",
            "Laboratorio 3 y tarea",
        ],
        "content": [
            {"title": "De datos en el código a datos en la nube", "kicker": "El salto de hoy", "bullets": [
                "Hasta ahora los datos vivían dentro del código",
                "Firebase Firestore los guarda en la nube",
                "Se organiza en colecciones y documentos",
                "La colección cafeterias: nombre, zona, categoria, especialidad, descripcion, imagenUrl",
            ]},
            {"title": "Tiempo real y CRUD", "kicker": "Cómo se comporta", "bullets": [
                "onSnapshot: la pantalla se actualiza sola",
                "No hace falta recargar la app",
                "CRUD: crear, leer, actualizar, borrar",
                "Hoy nos enfocamos en leer (lista) y crear (formulario)",
            ]},
            {"title": "Auto-provisioning de Firebase", "kicker": "Sin tocar la consola", "bullets": [
                "AI Studio ofrece \"Enable Firebase\" cuando lo necesitas",
                "Lo apruebas con un clic",
                "Configura Firebase automáticamente",
                "No entras a la consola de Firebase",
            ]},
            {"title": "Seguridad antes de publicar", "kicker": "Buenas prácticas", "bullets": [
                "No dejes claves en el código del cliente",
                "Usa Secrets en Settings para valores confidenciales",
                "Revisa las reglas de seguridad que generó la IA",
                "Profundizamos en esto en la sesión 6 (Investigación 1)",
            ]},
        ],
        "diagrams": [
            {"title": "Datos en tiempo real con Firestore", "kicker": "Cómo funciona",
             "steps": [("Firestore", "datos en la nube"), ("onSnapshot", "escucha cambios"),
                       ("La app", "se actualiza sola")],
             "footnote": "Cuando alguien crea o edita un dato, la pantalla se refresca sin recargar. "
                         "El contenido va en español; los campos en inglés (nombre, zona, categoria...)."},
            {"title": "Qué hace instalable a una PWA", "kicker": "Tres piezas + publicación",
             "steps": [("Manifest", "nombre, colores"), ("Íconos", "varios tamaños"),
                       ("Service Worker", "sin conexión"), ("Cloud Run", "URL para instalar")],
             "footnote": "Con esas tres piezas el teléfono ofrece \"Instalar\". Cloud Run publica tu app "
                         "con una dirección https propia."},
        ],
        "screens": [
            {"title": "Ruta del Café: con datos y lista para instalar", "kicker": "A dónde llegamos",
             "images": [("ruta-home", "Con datos y lista para instalar")],
             "bullets": [
                 "La lista sale de Firestore, no del código",
                 "El formulario (＋) crea registros nuevos",
                 "Aparecen sin recargar, gracias a onSnapshot",
                 "Manifest + íconos + service worker → instalable",
                 "Publicada en Cloud Run con su propia URL",
             ]},
        ],
        "pairs": [
            {"title": "Ruta del Café terminada: mock-up y app real",
             "mockup": "ruta-home", "mockup_caption": "Lo que diseñamos como referencia",
             "real": "ruta-real", "real_caption": "Con datos de Firestore y lista para instalar"},
        ],
        "demo": {
            "app": "Recetario Rápido",
            "image": "demo-recetario",
            "concept": [
                "Recetas guardadas en la nube con Firestore (PWA)",
                "Muestra datos en tiempo real y un formulario para crear registros",
                "El docente conecta Firebase en vivo con \"Enable Firebase\"",
            ],
        },
        "lab": {"title": "Laboratorio 3 — PWA con datos, instalable", "bullets": [
            "Guarda un checkpoint (cambio grande)",
            "Conecta tu app a Firestore y aprueba \"Enable Firebase\"",
            "Agrega el formulario para crear registros",
            "Verifica que aparece sin recargar (tiempo real)",
            "Hazla instalable como PWA y publícala en Cloud Run",
        ]},
        "closing": {
            "recap": [
                "Firestore guarda datos en la nube en tiempo real",
                "Manifest + íconos + service worker = PWA instalable",
                "El Proyecto 1 queda completo con esta sesión",
            ],
            "tareas": [
                "Tarea 3: define el modelo de datos de tu app",
                "Laboratorio 3: PWA con datos persistentes, instalable",
                "Examen Parcial (sesión 5): escrito en línea + entrega del Proyecto 1",
            ],
        },
    },
    {
        "session": 4,
        "title": "Desarrollo nativo I: fundamentos con IA",
        "agenda": [
            "Repaso del Proyecto 1",
            "PWA vs nativo",
            "Anatomía de una app Android",
            "App nativa en AI Studio",
            "El flujo de la vista previa",
            "Laboratorio 4 y tarea",
        ],
        "content": [
            {"title": "Anatomía de una app Android", "kicker": "En términos de diseño", "bullets": [
                "Kotlin + Jetpack Compose es la tecnología que usa la IA; tú diriges con prompts",
                "Material 3 es tu vocabulario de diseño: TopAppBar, cards, FAB, chips, sheets",
                "Pantallas y navegación; recursos: íconos, colores, tipografías",
                "No necesitas el código: necesitas el criterio de diseño",
            ]},
            {"title": "App nativa en AI Studio", "kicker": "Misma herramienta, del navegador", "bullets": [
                "Creas una app nueva y pides una app nativa de Android",
                "La IA la construye en Kotlin + Jetpack Compose (Material 3)",
                "Todo desde Chrome: no hay terminal, ni SDK, ni instalaciones",
                "Modelo por defecto: Gemini 3.5 Flash (Pro para tareas complejas)",
            ]},
            {"title": "El flujo de la vista previa", "kicker": "Cómo lo probamos", "bullets": [
                "Emulador de Android integrado en el navegador: se actualiza solo",
                "No creas ni arrancas nada; la vista previa refresca sola",
                "\"Install on Device\": instala en tu teléfono por USB (WebUSB)",
                "Sin ADB ni SDK; solo modo de desarrollador + depuración USB",
            ]},
            {"title": "Estructura de una app Android y su tema (Material 3)", "kicker": "Vocabulario de diseño, no programación", "bullets": [
                "El tema (theme) reúne colores, tipografías y formas: es el sistema visual de toda la app",
                "El color scheme de Material 3 —primario, secundario, fondo, superficie— es donde va tu paleta de marca",
                "Las pantallas (screens) y los componentes (Compose) se organizan en archivos separados",
                "Los recursos —íconos, colores, textos— viven aparte para poder reutilizarse",
                "Cambias el tema una vez y toda la app se ve consistente; lo diriges con prompts",
            ]},
        ],
        "compare": [
            {"title": "PWA vs app nativa", "kicker": "Cuándo cada una",
             "left_title": "PWA", "left": [
                 "Vive en el navegador",
                 "Se instala desde una URL",
                 "Alcance y velocidad de publicación",
                 "Ideal para catálogos y contenido",
                 "Una base de código para todos",
             ],
             "right_title": "Nativa (Android)", "right": [
                 "Se instala desde la tienda (Play Store)",
                 "Mejor acceso al teléfono: cámara, sensores",
                 "Presencia y descubrimiento en la tienda",
                 "Ideal para herramientas que capturan",
                 "Rendimiento y experiencia más ricos",
             ]},
        ],
        "screens": [
            {"title": "Bitácora Visual en el emulador del navegador", "kicker": "A dónde llegamos",
             "images": [("bitacora-grid", "Nativa, en el emulador del navegador")],
             "bullets": [
                 "Es una app nativa (Kotlin + Compose), pero tú la diriges con prompts",
                 "El emulador del navegador la muestra al instante",
                 "Cuadrícula de entradas con la estética de tu marca",
                 "Con \"Install on Device\" la instalas en tu teléfono por USB",
             ]},
        ],
        "demo": {
            "app": "Mi Estantería",
            "image": "demo-estanteria",
            "concept": [
                "Registro de libros leídos, app nativa de Android (Kotlin + Compose)",
                "Ilustra Material 3 como vocabulario de diseño y el emulador del navegador",
                "El docente la genera en vivo y la prueba con \"Install on Device\"",
            ],
        },
        "lab": {"title": "Laboratorio 4 — App nativa en el emulador del navegador", "bullets": [
            "Crea una app nueva en AI Studio (modo Build)",
            "Pide una app nativa de Android (Kotlin + Jetpack Compose) \"Bitácora Visual\"",
            "Revisa la bitácora de ejemplo en el emulador integrado del navegador",
            "Prueba \"Install on Device\" en tu teléfono por USB (opcional)",
            "Refina la cuadrícula de entradas y revisa con View diff",
        ]},
        "closing": {
            "recap": [
                "Nativo conviene para cámara, tienda y rendimiento",
                "La app nativa se construye en AI Studio, sin terminal ni SDK",
                "La IA construye nativo (Kotlin + Compose) igual que la PWA",
            ],
            "tareas": [
                "Tarea 4: investigación PWA vs nativo",
                "Laboratorio 4: app nativa en el emulador del navegador con su cuadrícula",
                "Examen Parcial en la próxima sesión: escrito en línea + entrega del Proyecto 1",
            ],
        },
    },
    {
        "session": 6,
        "title": "Desarrollo nativo II: producto real",
        "agenda": [
            "Retroalimentación del parcial",
            "Login con Google",
            "Datos por usuario en Firestore",
            "Fotos en Storage",
            "Reglas de seguridad",
            "Laboratorio 5 y tarea",
        ],
        "content": [
            {"title": "De bitácora de ejemplo a producto real", "kicker": "El salto de hoy", "bullets": [
                "Pasamos a un producto multiusuario",
                "Cada persona entra con su cuenta y ve solo su bitácora",
                "Datos en la nube y fotos reales",
                "El corazón del Proyecto Final",
            ]},
            {"title": "Login con Google (Authentication)", "kicker": "Identidad del usuario", "bullets": [
                "Pantalla de login con \"Iniciar sesión con Google\"",
                "Google Sign-In de Firebase Authentication",
                "Cada usuario obtiene un identificador único (uid)",
                "Detectar la sesión y permitir cerrar sesión",
            ]},
            {"title": "Datos separados por usuario", "kicker": "Firestore con uid", "bullets": [
                "Colección entradas: titulo, nota, imagenUrl, uid, fecha",
                "La bitácora filtra por uid: cada quien ve solo lo suyo",
                "Lectura en tiempo real",
                "Estado vacío amable para el usuario sin entradas",
            ]},
            {"title": "Reglas de seguridad", "kicker": "Proteger los datos", "bullets": [
                "Deciden quién puede leer y escribir qué",
                "Cada usuario solo accede a sus propias entradas",
                "Condición: uid == request.auth.uid",
                "Tema de tu Investigación 1: revísalas antes del Demo Day",
            ]},
        ],
        "diagrams": [
            {"title": "Datos y fotos privados por usuario", "kicker": "Cómo funciona",
             "steps": [("Login Google", "identidad"), ("uid", "quién eres"),
                       ("Firestore", "entradas del uid"), ("Storage", "las fotos")],
             "footnote": "La foto sube a Storage; en Firestore se guarda solo su dirección (imagenUrl) "
                         "junto al uid. Así cada quien ve únicamente sus propias entradas."},
        ],
        "screens": [
            {"title": "Login y bitácora privada", "kicker": "En pantalla",
             "images": [("bitacora-login", "Iniciar sesión con Google"),
                        ("bitacora-grid", "Solo mis entradas (filtradas por uid)")]},
        ],
        "pairs": [
            {"title": "Bitácora Visual: la meta y el resultado",
             "mockup": "bitacora-grid", "mockup_caption": "El mock-up de referencia (la guía)",
             "real": "bitacora-real", "real_caption": "La app nativa generada por la IA"},
        ],
        "demo": {
            "app": "Diario de Viajes",
            "image": "demo-diario",
            "concept": [
                "Diario privado por usuario, app nativa de Android",
                "Ilustra login con Google y datos separados por uid en Firestore",
                "El docente entra con dos cuentas para mostrar que cada quien ve solo lo suyo",
            ],
        },
        "lab": {"title": "Laboratorio 5 — App nativa con datos y login", "bullets": [
            "Agrega login con Google y aprueba la activación de Firebase",
            "Conecta las entradas a Firestore filtrando por uid",
            "Agrega el formulario con subida de foto a Storage",
            "Verifica que la entrada nueva aparece automáticamente",
            "Prueba con dos cuentas que cada quien ve solo lo suyo",
        ]},
        "closing": {
            "recap": [
                "Login con Google y uid separan los datos por usuario",
                "Storage guarda fotos; Firestore guarda su dirección",
                "Las reglas de seguridad protegen a cada usuario",
            ],
            "tareas": [
                "Investigación 1: reglas de seguridad de Firebase",
                "Laboratorio 5: app nativa con datos y login",
                "Quiz de la unidad nativa con datos",
            ],
        },
    },
    {
        "session": 7,
        "title": "Desarrollo nativo III: pulido y capacidades del teléfono",
        "agenda": [
            "Repaso de la sesión 6",
            "Capacidades del teléfono",
            "Cámara y permisos",
            "Depurar con evidencia",
            "Pulido de marca",
            "Laboratorio 6 y cierre del desarrollo",
        ],
        "content": [
            {"title": "Capacidades del teléfono", "kicker": "Lo que solo lo nativo da bien", "bullets": [
                "Cámara, notificaciones, ubicación",
                "Hoy nos enfocamos en la cámara",
                "Es la diferencia visible frente a una web",
                "Mejora la experiencia y el valor del producto",
            ]},
            {"title": "Cámara y permisos", "kicker": "UX de confianza", "bullets": [
                "Las funciones sensibles requieren pedir permiso",
                "Mensaje claro en español: para qué se usa",
                "La foto tomada se sube a Storage igual que las demás",
                "Si el usuario niega, deja disponible elegir desde la galería",
            ]},
            {"title": "Depurar con evidencia, no a ciegas", "kicker": "Herramientas", "bullets": [
                "Toma capturas de la vista previa del navegador (o del teléfono)",
                "Modo de anotación: haz clic sobre el elemento y describe el ajuste",
                "Le pasas esa evidencia a la IA para corregir con precisión",
                "Apóyate también en View diff y la pestaña Code",
            ]},
            {"title": "Pulido de marca", "kicker": "De ejercicio a producto", "bullets": [
                "Ícono de app propio con tu color de marca",
                "Splash screen al abrir la app",
                "Color de marca consistente en barras y botones",
                "Puedes generar el ícono con un AI Chip (Nano Banana)",
            ]},
        ],
        "diagrams": [
            {"title": "De la cámara a la bitácora", "kicker": "El recorrido de una foto",
             "steps": [("Cámara", "pide permiso"), ("Foto", "tomada"),
                       ("Storage", "sube el archivo"), ("Bitácora", "aparece sola")],
             "footnote": "En Firestore se guarda la dirección de la foto (imagenUrl); la imagen vive en "
                         "Storage. La entrada nueva aparece en la cuadrícula sin recargar."},
        ],
        "screens": [
            {"title": "Capturar con la cámara", "kicker": "En pantalla",
             "images": [("bitacora-captura", "Nueva entrada: foto + nota")],
             "bullets": [
                 "La cámara pide permiso con un mensaje claro en español",
                 "La foto tomada encabeza la nueva entrada",
                 "Título y nota acompañan la imagen",
                 "Al guardar, sube a Storage y aparece en la bitácora",
                 "Ícono y splash con tu marca la vuelven un producto",
             ]},
        ],
        "pairs": [
            {"title": "Bitácora Visual pulida: mock-up y app real",
             "mockup": "bitacora-captura", "mockup_caption": "El mock-up de captura (referencia)",
             "real": "bitacora-real", "real_caption": "La app nativa pulida, ya un producto"},
        ],
        "demo": {
            "app": "Foto del Día",
            "image": "demo-foto",
            "concept": [
                "Una foto al día, app nativa cámara-first de Android",
                "Ilustra permisos de cámara, subida a Storage y pulido de marca",
                "El docente toma una foto en vivo y la ve aparecer en la app",
            ],
        },
        "lab": {"title": "Laboratorio 6 — App nativa con cámara, pulida", "bullets": [
            "Agrega \"Tomar foto\" con la cámara (con permiso)",
            "Verifica que la foto tomada sube a Storage",
            "Diagnostica con capturas de la vista previa y el modo de anotación",
            "Pule el ícono y el splash con tu marca",
            "Toma una captura final para tu portafolio",
        ]},
        "closing": {
            "recap": [
                "La cámara requiere permiso con un mensaje claro",
                "Depura con capturas de la vista previa y anotación, no a ciegas",
                "Ícono y splash convierten el ejercicio en producto",
            ],
            "tareas": [
                "Laboratorio 6: app nativa pulida con una capacidad del teléfono",
                "Reúne capturas y notas para tu portafolio",
                "Sesión 10: Examen Final escrito en línea + Demo Day (entrega del Proyecto Final)",
            ],
        },
    },
    {
        "session": 8,
        "title": "Distribución y marketing de apps",
        "agenda": [
            "Marco de la sesión",
            "El tamaño del mercado de apps",
            "PWA vs Play Store",
            "Store listing y ASO básico",
            "Casos reales",
            "Mini-taller: ficha de tienda",
        ],
        "content": [
            {"title": "Distribuir y comunicar el producto", "kicker": "Marco", "bullets": [
                "Ya construyeron el producto; ahora toca distribuirlo",
                "Tema más ligero pero clave para tu empleabilidad",
                "El marketing de apps es comunicación visual: tu terreno",
            ]},
            {"title": "Store listing", "kicker": "El empaque de la app", "bullets": [
                "Título, descripción corta y descripción larga",
                "Ícono y capturas que comunican valor",
                "Categoría adecuada en la tienda",
                "Es lo primero que ve quien decide instalar",
            ]},
            {"title": "ASO básico", "kicker": "Que te encuentren", "bullets": [
                "Palabras clave en título y descripción",
                "Capturas que comunican valor en 3 segundos",
                "Ícono distintivo y reconocible",
                "Habla de beneficios, no solo de funciones",
            ]},
        ],
        "stats": [
            {"title": "El mercado de las apps", "kicker": "El tamaño del terreno", "cards": [
                {"num": "2.26 M", "label": "apps disponibles en Google Play, la tienda más grande",
                 "source": "Statista, Google Play (2024)"},
                {"num": "136 mil M", "label": "descargas de apps en el mundo en un año",
                 "source": "Sensor Tower, State of Mobile 2025"},
                {"num": "5 h/día", "label": "en promedio pasamos en el teléfono cada día",
                 "source": "data.ai / Sensor Tower, State of Mobile 2024"},
            ]},
        ],
        "compare": [
            {"title": "Dos vías de distribución", "kicker": "PWA vs Play Store",
             "left_title": "PWA", "left": [
                 "Se comparte por URL",
                 "Se instala desde el navegador",
                 "Publicación inmediata, sin revisión",
                 "Tu PWA ya es distribuible hoy",
             ],
             "right_title": "Play Store", "right": [
                 "Tienda oficial de Android",
                 "Requiere cuenta y revisión",
                 "Presencia y descubrimiento",
                 "La nativa apunta a la tienda",
             ]},
        ],
        "lab": {"title": "Mini-taller — Ficha de tienda (store listing)", "bullets": [
            "Título con una palabra clave",
            "Descripción corta (1 frase de gancho) y descripción larga",
            "Ícono ya pulido en la sesión 7",
            "2–3 capturas que comuniquen el valor (salen de la vista previa o del teléfono)",
            "Categoría sugerida en la tienda",
        ]},
        "closing": {
            "recap": [
                "PWA y Play Store: dos vías con usos distintos",
                "Store listing y ASO comunican el valor de tu app",
                "El marketing de apps es comunicación y diseño",
            ],
            "tareas": [
                "Tarea 5: estrategia de distribución y marketing",
                "Termina tu store listing (alimenta el portafolio)",
                "Quiz de la unidad de distribución y portafolio",
            ],
        },
    },
    {
        "session": 9,
        "title": "Portafolio, pitch y empleabilidad",
        "agenda": [
            "Marco de la sesión",
            "El portafolio del diseñador-creador con IA",
            "Ética y accesibilidad",
            "Anatomía de un pitch",
            "Empleabilidad: los datos",
            "Mini-taller: portafolio + pitch",
        ],
        "content": [
            {"title": "Convierte tu proyecto en una pieza de portafolio", "kicker": "Marco", "bullets": [
                "Tu app es también tu portafolio: doble valor",
                "Cuenta problema → proceso → resultado",
                "No solo pantallas bonitas: muestra tu proceso",
            ]},
            {"title": "Tu diferenciador", "kicker": "Lo que te hace destacar", "bullets": [
                "Sabes dirigir a la IA para construir productos reales",
                "Muestra tu bitácora de prompts y decisiones de diseño",
                "Es desarrollo de software asistido por IA, no \"vibe coding\"",
            ]},
            {"title": "Ética y accesibilidad", "kicker": "Profesionalismo", "bullets": [
                "Contraste, tamaños táctiles y lenguaje claro",
                "Uso responsable de la IA y de imágenes generadas",
                "Parte de presentarte como profesional serio",
            ]},
            {"title": "Anatomía de un pitch (2–3 min)", "kicker": "Cómo presentar", "bullets": [
                "Gancho: qué es y para quién, en una frase",
                "Demo en vivo: login, bitácora, agregar entrada con cámara",
                "Proceso: cómo lo dirigiste con IA, una decisión clave",
                "Cierre: por qué representa tu trabajo creativo",
            ]},
        ],
        "stats": [
            {"title": "Por qué esta habilidad te posiciona", "kicker": "Empleabilidad", "cards": [
                {"num": "+38%", "label": "crecieron las ofertas de empleo con IA (2020–2024)",
                 "source": "LinkedIn, Future of Work: AI at Work 2025"},
                {"num": "59%", "label": "de diseñadores y desarrolladores ya usa IA en su trabajo",
                 "source": "Figma, AI Design Report 2024"},
                {"num": "50 ms", "label": "tarda un usuario en formar su primera impresión visual",
                 "source": "Lindgaard et al., 2006"},
            ]},
        ],
        "lab": {"title": "Mini-taller — Caso de portafolio + pitch", "bullets": [
            "Caso: problema, proceso (con prompts y capturas) y resultado",
            "Guion de pitch de 2–3 minutos (gancho, demo, proceso, cierre)",
            "Ensaya el pitch en voz alta",
            "Recibe retroalimentación del grupo (¿se entiende el valor?)",
            "Es el ensayo general del Demo Day",
        ]},
        "closing": {
            "recap": [
                "Un buen caso cuenta problema, proceso y resultado",
                "El pitch comunica valor en 2–3 minutos",
                "Tu portafolio es una pieza real para tu carrera",
            ],
            "tareas": [
                "Termina tu caso de portafolio y tu guion de pitch",
                "Ensaya con emulador o teléfono listo",
                "Sesión 10: Examen Final escrito en línea + Demo Day (entrega + pitch)",
            ],
        },
    },
]


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    created = []
    for spec in DECKS:
        out_path = os.path.join(out_dir, f"sesion-{spec['session']:02d}.pptx")
        build_deck(spec, out_path)
        created.append(out_path)
        print(f"Creado: {out_path}")
    print(f"\nTotal de decks generados: {len(created)}")


if __name__ == "__main__":
    main()
